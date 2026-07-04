from pathlib import Path
from typing import Iterable, List, Tuple


TEXT_OUTPUT_FORMATS = ("markdown", "ppt-outline", "speaker-notes", "qa")
BINARY_OUTPUT_FORMATS = ("docx", "pdf", "pptx")


def write_output(path: Path, content: str, output_format: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    if output_format in TEXT_OUTPUT_FORMATS:
        path.write_text(content, encoding="utf-8")
        return
    if output_format == "docx":
        write_docx(path, content)
        return
    if output_format == "pdf":
        write_pdf(path, content)
        return
    if output_format == "pptx":
        write_pptx(path, content)
        return
    raise ValueError("Unsupported output format: {0}".format(output_format))


def render_format_for_output(output_format: str) -> str:
    if output_format == "pptx":
        return "ppt-outline"
    if output_format in ("docx", "pdf"):
        return "markdown"
    return output_format


def write_docx(path: Path, content: str) -> None:
    try:
        from docx import Document
    except ImportError as error:
        raise ImportError("DOCX export requires python-docx. Install it with: pip3 install python-docx") from error

    document = Document()
    for line in content.splitlines():
        if line.startswith("# "):
            document.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            document.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            document.add_heading(line[4:].strip(), level=3)
        elif line.startswith("- "):
            document.add_paragraph(line[2:].strip(), style="List Bullet")
        elif line.strip():
            document.add_paragraph(line.strip())
    document.save(str(path))


def write_pdf(path: Path, content: str) -> None:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError as error:
        raise ImportError("PDF export requires reportlab. Install it with: pip3 install reportlab") from error

    pdf = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    x = 54
    y = height - 54
    line_height = 14

    for line in _wrap_lines(content, max_chars=92):
        if y < 54:
            pdf.showPage()
            y = height - 54
        pdf.drawString(x, y, line)
        y -= line_height
    pdf.save()


def write_pptx(path: Path, content: str) -> None:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise ImportError("PPTX export requires python-pptx. Install it with: pip3 install python-pptx") from error

    presentation = Presentation()
    for title, bullets in _slides_from_outline(content):
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        if not bullets:
            body.text = ""
            continue
        body.text = bullets[0]
        for bullet in bullets[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
    presentation.save(str(path))


def _wrap_lines(content: str, max_chars: int) -> Iterable[str]:
    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            yield ""
            continue
        while len(line) > max_chars:
            split_at = line.rfind(" ", 0, max_chars)
            if split_at <= 0:
                split_at = max_chars
            yield line[:split_at]
            line = line[split_at:].strip()
        yield line


def _slides_from_outline(content: str) -> List[Tuple[str, List[str]]]:
    slides = []
    current_title = None
    current_bullets = []
    for line in content.splitlines():
        if line.startswith("## Slide"):
            if current_title:
                slides.append((current_title, current_bullets))
            current_title = line[3:].strip()
            current_bullets = []
        elif line.startswith("# ") and current_title is None:
            current_title = line[2:].strip()
            current_bullets = []
        elif line.startswith("- ") and current_title:
            current_bullets.append(line[2:].strip())
    if current_title:
        slides.append((current_title, current_bullets))
    return slides or [("Research Brief", [line for line in content.splitlines() if line.strip()])]
