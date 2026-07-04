import csv
import html
import re
from pathlib import Path
from typing import Iterable, List
from urllib.parse import urlparse
from urllib.request import Request, urlopen

from .models import Evidence, Source

MIN_EVIDENCE_LENGTH = 12
MAX_LABEL_LENGTH = 80


def read_source_notes(paths: Iterable[object]) -> List[Evidence]:
    evidence = []
    for index, raw_path in enumerate(paths, start=1):
        path_text = str(raw_path)
        if _is_url(path_text):
            evidence.extend(_read_url(path_text, index))
            continue

        path = Path(path_text)
        if not path.exists():
            raise FileNotFoundError("Source file not found: {0}".format(path))
        if not path.is_file():
            raise ValueError("Source path is not a file: {0}".format(path))
        source = Source(
            id="source-{0}".format(index),
            title=path.name,
            path=str(path),
            source_type=_source_type_for_path(path),
        )
        evidence.extend(_read_file_evidence(path, source))
    return evidence


def _read_file_evidence(path: Path, source: Source) -> List[Evidence]:
    suffix = path.suffix.lower()
    if suffix in (".txt", ".md", ".markdown", ".rst"):
        return _read_text_evidence(path, source)
    if suffix == ".csv":
        return _read_csv_evidence(path, source)
    if suffix == ".pdf":
        return _read_pdf_evidence(path, source)
    if suffix == ".docx":
        return _read_docx_evidence(path, source)
    if suffix == ".pptx":
        return _read_pptx_evidence(path, source)
    if suffix in (".xlsx", ".xlsm"):
        return _read_xlsx_evidence(path, source)
    raise ValueError("Unsupported source file type: {0}".format(path.suffix or path))


def _read_text_evidence(path: Path, source: Source) -> List[Evidence]:
    content = path.read_text(encoding="utf-8")
    return _evidence_from_text(content, source, location_prefix="line")


def _evidence_from_text(
    content: str,
    source: Source,
    location_prefix: str = "line",
) -> List[Evidence]:
    items = []
    in_code_block = False
    paragraph_lines = []
    paragraph_start = 0

    for line_number, raw_line in enumerate(content.splitlines(), start=1):
        text = " ".join(raw_line.strip("-* \t").split())
        if _is_code_fence(text):
            _append_paragraph(items, paragraph_lines, paragraph_start, source, location_prefix)
            paragraph_lines = []
            paragraph_start = 0
            in_code_block = not in_code_block
            continue
        if in_code_block:
            continue
        if not _is_evidence_line(text):
            _append_paragraph(items, paragraph_lines, paragraph_start, source, location_prefix)
            paragraph_lines = []
            paragraph_start = 0
            continue

        if _is_bullet_line(raw_line):
            _append_paragraph(items, paragraph_lines, paragraph_start, source, location_prefix)
            paragraph_lines = []
            paragraph_start = 0
            _append_evidence(items, text, line_number, line_number, source, location_prefix)
            continue

        if not paragraph_lines:
            paragraph_start = line_number
        paragraph_lines.append(text)

    _append_paragraph(items, paragraph_lines, paragraph_start, source, location_prefix)
    return items


def _read_csv_evidence(path: Path, source: Source) -> List[Evidence]:
    items = []
    with path.open("r", encoding="utf-8-sig", newline="") as file:
        reader = csv.reader(file)
        headers = []
        for row_number, row in enumerate(reader, start=1):
            values = [cell.strip() for cell in row if cell.strip()]
            if not values:
                continue
            if row_number == 1:
                headers = values
                continue
            text = _row_to_text(headers, values)
            if _is_evidence_line(text):
                _append_custom_evidence(items, text, "row {0}".format(row_number), source)
    return items


def _read_pdf_evidence(path: Path, source: Source) -> List[Evidence]:
    try:
        from pypdf import PdfReader
    except ImportError as error:
        raise ImportError("PDF ingestion requires pypdf. Install it with: pip3 install pypdf") from error

    items = []
    reader = PdfReader(str(path))
    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        for index, paragraph in enumerate(_split_paragraphs(text), start=1):
            if _is_evidence_line(paragraph):
                _append_custom_evidence(
                    items,
                    paragraph,
                    "page {0}, paragraph {1}".format(page_index, index),
                    source,
                )
    return items


def _read_docx_evidence(path: Path, source: Source) -> List[Evidence]:
    try:
        from docx import Document
    except ImportError as error:
        raise ImportError("DOCX ingestion requires python-docx. Install it with: pip3 install python-docx") from error

    items = []
    document = Document(str(path))
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = " ".join(paragraph.text.split())
        if _is_evidence_line(text):
            _append_custom_evidence(items, text, "paragraph {0}".format(index), source)
    return items


def _read_pptx_evidence(path: Path, source: Source) -> List[Evidence]:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise ImportError("PPTX ingestion requires python-pptx. Install it with: pip3 install python-pptx") from error

    items = []
    presentation = Presentation(str(path))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(" ".join(shape.text.split()))
        text = " ".join(texts)
        if _is_evidence_line(text):
            _append_custom_evidence(items, text, "slide {0}".format(slide_index), source)
    return items


def _read_xlsx_evidence(path: Path, source: Source) -> List[Evidence]:
    try:
        from openpyxl import load_workbook
    except ImportError as error:
        raise ImportError("XLSX ingestion requires openpyxl. Install it with: pip3 install openpyxl") from error

    items = []
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        for sheet in workbook.worksheets:
            rows = sheet.iter_rows(values_only=True)
            headers = []
            for row_number, row in enumerate(rows, start=1):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if not values:
                    continue
                if row_number == 1:
                    headers = values
                    continue
                text = _row_to_text(headers, values)
                if _is_evidence_line(text):
                    _append_custom_evidence(
                        items,
                        text,
                        "{0}!row {1}".format(sheet.title, row_number),
                        source,
                    )
    finally:
        workbook.close()
    return items


def _read_url(url: str, index: int) -> List[Evidence]:
    request = Request(url, headers={"User-Agent": "ResearchBriefingAgent/0.4"})
    with urlopen(request, timeout=30) as response:
        content_type = response.headers.get("Content-Type", "")
        raw = response.read().decode(_charset_from_content_type(content_type), errors="replace")

    source = Source(
        id="source-{0}".format(index),
        title=urlparse(url).netloc or url,
        path=url,
        source_type="url",
    )
    text = _html_to_text(raw) if "html" in content_type.lower() else raw
    return _evidence_from_text(text, source, location_prefix="line")


def _is_evidence_line(text: str) -> bool:
    if not text:
        return False
    if text.startswith("#"):
        return False
    if len(text) < MIN_EVIDENCE_LENGTH:
        return False
    if text.endswith(":") and len(text) <= MAX_LABEL_LENGTH:
        return False
    return True


def _is_code_fence(text: str) -> bool:
    return text.startswith("```") or text.startswith("~~~")


def _is_bullet_line(raw_line: str) -> bool:
    stripped = raw_line.lstrip()
    return stripped.startswith("- ") or stripped.startswith("* ")


def _append_paragraph(
    items: List[Evidence],
    paragraph_lines: List[str],
    start_line: int,
    source: Source,
    location_prefix: str,
) -> None:
    if not paragraph_lines:
        return
    text = " ".join(paragraph_lines)
    _append_evidence(
        items,
        text,
        start_line,
        start_line + len(paragraph_lines) - 1,
        source,
        location_prefix,
    )


def _append_evidence(
    items: List[Evidence],
    text: str,
    start_line: int,
    end_line: int,
    source: Source,
    location_prefix: str,
) -> None:
    location = _format_location(start_line, end_line, location_prefix)
    _append_custom_evidence(items, text, location, source)


def _append_custom_evidence(
    items: List[Evidence],
    text: str,
    location: str,
    source: Source,
) -> None:
    evidence_id = "{0}-{1}".format(source.id, _slug(location))
    items.append(
        Evidence(
            id=evidence_id,
            text=text,
            source=source,
            location=location,
            confidence="source-backed",
        )
    )


def _format_location(start_line: int, end_line: int, prefix: str) -> str:
    if start_line == end_line:
        return "{0} {1}".format(prefix, start_line)
    if prefix == "line":
        return "lines {0}-{1}".format(start_line, end_line)
    return "{0}s {1}-{2}".format(prefix, start_line, end_line)


def _source_type_for_path(path: Path) -> str:
    suffix = path.suffix.lower().lstrip(".")
    return suffix or "file"


def _is_url(value: str) -> bool:
    parsed = urlparse(value)
    return parsed.scheme in ("http", "https") and bool(parsed.netloc)


def _row_to_text(headers: List[str], values: List[str]) -> str:
    if headers and len(headers) == len(values):
        return "; ".join(
            "{0}: {1}".format(header, value)
            for header, value in zip(headers, values)
        )
    return "; ".join(values)


def _split_paragraphs(text: str) -> List[str]:
    chunks = re.split(r"\n\s*\n|(?<=\.)\s{2,}", text)
    return [" ".join(chunk.split()) for chunk in chunks if chunk.strip()]


def _charset_from_content_type(content_type: str) -> str:
    match = re.search(r"charset=([^;]+)", content_type, flags=re.IGNORECASE)
    if match:
        return match.group(1).strip()
    return "utf-8"


def _html_to_text(value: str) -> str:
    value = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", value)
    value = re.sub(r"(?i)<br\s*/?>", "\n", value)
    value = re.sub(r"(?i)</p>|</div>|</li>|</h[1-6]>", "\n", value)
    value = re.sub(r"<[^>]+>", " ", value)
    return html.unescape(value)


def _slug(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower() or "item"
